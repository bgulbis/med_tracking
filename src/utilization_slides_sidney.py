from datetime import date
import glob
import numpy as np
import pandas as pd

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

filepaths = glob.glob("data/tidy/sidney/med-utilization_monthly_sidney_*.csv")
df = pd.concat(map(lambda x: pd.read_csv(x, index_col=0, parse_dates=True), filepaths), sort=False)
df.sort_index(inplace=True)
df['MONTH'] = pd.to_datetime('2018-' + (df.index-pd.DateOffset(months=6)).strftime('%m-%d')) + pd.DateOffset(months=6)
df['FY'] = 'FY' + (df.index + pd.DateOffset(months=6)).strftime('%y')

when = df.index[-1]
date_cut = date(day=1, month=7, year=(when+pd.DateOffset(months=6)).year) - pd.DateOffset(years=4)
df = df.loc[date_cut:]

df['MEDICATION'] = df['MEDICATION'].str.title()
df['MEDICATION'] = df['MEDICATION'].str.replace('Acetaminophen', 'Acetaminophen IV')
df['MEDICATION'] = df['MEDICATION'].str.replace('Levothyroxine', 'Levothyroxine IV')
df['MEDICATION'] = df['MEDICATION'].str.replace(' Human', '')
df['MEDICATION'] = df['MEDICATION'].str.replace('Bupivacaine Liposome', 'Bupivacaine (liposomal)')
df['MEDICATION'] = df['MEDICATION'].str.replace('Immune Globulin Intravenous And Subcut', 'IVIG')
df['MEDICATION'] = df['MEDICATION'].str.replace('Immune Globulin Intravenous', 'IVIG')

df['INPT'] = df['ENCOUNTER_TYPE'].isin(['Inpatient', 'Observation', 'Emergency'])
df.loc[(df['INPT'] == True) & (df['MEDICATION'] == "IVIG"), 'MEDICATION'] = 'IVIG (inpatient)'
df.loc[(df['INPT'] == False) & (df['MEDICATION'] == "IVIG"), 'MEDICATION'] = 'IVIG (outpatient)'

meds = np.sort(df['MEDICATION'].unique())

df_pvt = pd.pivot_table(data=df, 
                        values='DOSES', 
                        index=['MEDICATION', 
                               pd.Grouper(key='MONTH', freq='MS')], 
                        columns=pd.Grouper(key='FY'), 
                        aggfunc=np.sum, 
                        fill_value=0,
                        dropna=False)

prs = Presentation("doc/template.pptx")
# title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Utilization of Target Medications"
data_end = df.index[-1].strftime('%B %Y')
subtitle.text = "Data through: " + data_end + "\nBrian Gulbis, PharmD, BCPS"

# blank slide layout
blank_slide_layout = prs.slide_layouts[6]
# chart size
x, y, cx, cy = Inches(0.5), Inches(1), Inches(9), Inches(6)

for i in meds:
    df_i = df_pvt.loc[i]

    if np.sum(df_i.sum()) == 0:
        continue

    # add new slide for graph
    slide = prs.slides.add_slide(blank_slide_layout)

    chart_data = CategoryChartData()
    chart_data.categories = df_i.index

    for j in range(0, len(df_i.columns)):
        if j == 3:
            ser = df_i.iloc[:, j]
            idx = pd.to_datetime('2018-' + (when-pd.DateOffset(months=6)).strftime('%m-%d')) + pd.DateOffset(months=6)
            idx_next = idx + pd.DateOffset(months=1)
            keep = ser.loc[:idx]
            chg = ser.loc[idx_next:].replace({0: None})
            keep = keep.append(chg)
            chart_data.add_series(df_i.columns[j], keep)
        else:
            chart_data.add_series(df_i.columns[j], df_i.iloc[:, j])

    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data).chart
    chart.chart_title.text_frame.text = "{0} utilization".format(i)

prs.save("report/sidney/utilization_slides.pptx")
