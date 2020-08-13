# %%
from datetime import date, datetime
import glob
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from rpy2.robjects import pandas2ri
from rpy2.robjects.packages import importr
from rpy2.robjects.functions import SignatureTranslatedFunction

from statsmodels.tsa.arima_model import ARIMA, ARIMAResults
from pmdarima.arima import AutoARIMA
from pmdarima.pipeline import Pipeline
from pmdarima.preprocessing import BoxCoxEndogTransformer

# %%
pandas2ri.activate()
forecast = importr('forecast')
base = importr('base')
forecast.auto_arima = SignatureTranslatedFunction(forecast.auto_arima, init_prm_translate = {'lambda_': 'lambda'})
forecast.BoxCox = SignatureTranslatedFunction(forecast.BoxCox, init_prm_translate = {'lambda_': 'lambda'})

# %%
# import data
filepaths = glob.glob("../data/raw/tmc_target_meds_*.csv")
df = pd.concat(map(lambda x: pd.read_csv(x, parse_dates=['DOSE_DATETIME']), filepaths), sort=False)

# %%
df['MEDICATION'] = df['MEDICATION'].str.title()
df['MEDICATION'] = df['MEDICATION'].str.replace('Acetaminophen', 'Acetaminophen IV')
df['MEDICATION'] = df['MEDICATION'].str.replace('Levothyroxine', 'Levothyroxine IV')
df['MEDICATION'] = df['MEDICATION'].str.replace(' Human', '')
df['MEDICATION'] = df['MEDICATION'].str.replace('Bupivacaine Liposome', 'Bupivacaine (liposomal)')
df['MEDICATION'] = df['MEDICATION'].str.replace('Immune Globulin Intravenous And Subcut', 'IVIG')
df['MEDICATION'] = df['MEDICATION'].str.replace('Immune Globulin Intravenous', 'IVIG')

df['INPT'] = df['ENCNTR_TYPE'].isin(['Inpatient', 'Observation', 'Emergency'])
df.loc[(df['INPT'] == True) & (df['MEDICATION'] == "IVIG"), 'MEDICATION'] = 'IVIG (inpatient)'
df.loc[(df['INPT'] == False) & (df['MEDICATION'] == "IVIG"), 'MEDICATION'] = 'IVIG (outpatient)'

df = df.sort_values(['MEDICATION', 'DOSE_DATETIME'])

meds = np.sort(df['MEDICATION'].unique())

# %%
df_doses = df.copy()
df_doses = df_doses.groupby('MEDICATION').resample('MS', on='DOSE_DATETIME').count()
df_doses = df_doses[['EVENT_ID']].rename(columns={'EVENT_ID': 'DOSES'})
# df_doses = df_doses.reset_index(level='DOSE_DATETIME')
# df_doses['MONTH'] = pd.to_datetime('2020-' + (df_doses['DOSE_DATETIME'] - pd.DateOffset(months=6)).dt.strftime('%m-%d')) + pd.DateOffset(months=6)
# df_doses['FY'] = 'FY' + (df_doses['DOSE_DATETIME'] + pd.DateOffset(months=6)).dt.strftime('%y')

when = df['DOSE_DATETIME'].max()
date_cut = date(day=1, month=7, year=(when + pd.DateOffset(months=6)).year) - pd.DateOffset(years=4)
# df_doses = df_doses[df_doses['DOSE_DATETIME'] >= date_cut]

# %%
i = 'Sugammadex'
df_actual = df_doses.loc[i].copy()
df_actual.columns = ["Actual"]

if df_actual.index[-1] != df_doses.index[-1][1]:
    new_idx = pd.date_range(df_actual.index[0], when, freq='MS')
    df_actual = df_actual.reindex(new_idx, fill_value=0)

# %%
# R forecast

model = forecast.auto_arima(
    df_actual["Actual"],
    max_order=10,
    stepwise=False,
    approximation=False,
    lambda_="auto",
    biasadj=True
)

fcast = forecast.forecast(model, h=12)

idx = pd.date_range(df_actual.index[-1] + pd.DateOffset(months=1), periods=12, freq='MS')
yhat = pd.Series(fcast.rx2("mean"), name="Forecast", index=idx)
lwr = pd.Series(fcast.rx2("lower")[0:12], name="Lower", index=idx)
upr = pd.Series(fcast.rx2("upper")[0:12], name="Upper", index=idx)
# lwr = pd.Series(fcast.rx2("lower")[:, :1].reshape(12, ), name="Lower", index=idx)
# upr = pd.Series(fcast.rx2("upper")[:, :1].reshape(12, ), name="Upper", index=idx)
df_fcast = pd.concat([yhat, lwr, upr], axis=1, sort=False)

df_combined = pd.concat([df_actual, df_fcast], axis=1, sort=True)
# df_combined = df_combined.replace({pd.np.nan: None})
df_combined = df_combined.loc[date_cut:]

# %%
# pmdarima

bc = BoxCoxEndogTransformer(lmbda2=1, neg_action='ignore')
bc.fit(df_actual)

pipeline = Pipeline([
    ('boxcox', BoxCoxEndogTransformer(lmbda2=1, neg_action='ignore')),
    ('arima', AutoARIMA(start_p=1, start_q=1, m=12, max_order=10, stepwise=False))
])

# mod = AutoARIMA(start_p=1, start_q=1, m=12, max_order=10)
# mod.fit(df_bc)
pipeline.fit(df_actual)
preds, conf_int = pipeline.predict(n_periods=12, return_conf_int=True, alpha=0.2)

idx = pd.date_range(df_actual.index[-1] + pd.DateOffset(months=1), periods=12, freq='MS')
preds = pd.Series(bc.inverse_transform(preds)[0], name="Forecast", index=idx)
conf_low = pd.Series(bc.inverse_transform(conf_int[:, 0])[0], name="Lower", index=idx)
conf_high = pd.Series(bc.inverse_transform(conf_int[:, 1])[0], name="Upper", index=idx)

df_fcast_py = pd.concat([preds, conf_low, conf_high], axis=1, sort=False)

df_combined_py = pd.concat([df_actual, df_fcast_py], axis=1, sort=True)
# df_combined = df_combined.replace({pd.np.nan: None})
df_combined_py = df_combined_py.loc[date_cut:]

# %%
plt.plot(df_combined)

# %%
plt.plot(df_combined_py)

# %%
for i in meds:
    df_actual = df_doses.loc[i].copy()

    # if df_actual['DOSES'].sum() == 0:
    #     continue

    # df_actual = df_actual.resample('MS', on='DOSE_DATETIME').count()
    # df_actual = df_actual.drop(columns=['PATIENTS', 'INPT'])
    df_actual.columns = ["Actual"]

    # if len(df_actual.index) < len(df.index):
    if df_actual.index[-1] != df_doses.index[-1][1]:
        new_idx = pd.date_range(df_actual.index[0], when, freq='MS')
        df_actual = df_actual.reindex(new_idx, fill_value=0)

    # pipeline = Pipeline([
    #     ('boxcox', BoxCoxEndogTransformer()),
    #     ('arima', AutoARIMA(start_p=1, start_q=1, m=12, max_order=10))
    # ])

    pipeline = AutoARIMA(start_p=1, start_q=1, m=12, max_order=10)
    pipeline.fit(df_actual)
    preds, conf_int = pipeline.predict(n_periods=12, return_conf_int=True)

    model = forecast.auto_arima(
        df_actual["Actual"],
        max_order=10,
        stepwise=False,
        approximation=False,
        lambda_="auto",
        biasadj=True
    )

    fcast = forecast.forecast(model, h=12)
    
    idx = pd.date_range(df_actual.index[-1] + pd.DateOffset(months=1), periods=12, freq='MS')
    yhat = pd.Series(fcast.rx2("mean"), name="Forecast", index=idx)
    lwr = pd.Series(fcast.rx2("lower")[0:12], name="Lower", index=idx)
    upr = pd.Series(fcast.rx2("upper")[0:12], name="Upper", index=idx)
    # lwr = pd.Series(fcast.rx2("lower")[:, :1].reshape(12, ), name="Lower", index=idx)
    # upr = pd.Series(fcast.rx2("upper")[:, :1].reshape(12, ), name="Upper", index=idx)
    df_fcast = pd.concat([yhat, lwr, upr], axis=1, sort=False)

    # if len(df_actual.index) < len(df.index):
    if df_actual.index[0] != df_doses.index[0][1]:
        new_idx = pd.date_range(date_cut, when, freq='MS')
        df_actual = df_actual.reindex(new_idx, fill_value=0)

    df_combined = pd.concat([df_actual, df_fcast], axis=1, sort=True)
    df_combined = df_combined.replace({pd.np.nan: None})
    df_combined = df_combined.loc[date_cut:]
        
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chart_data = CategoryChartData()
    chart_data.categories = df_combined.index

    chart_data.add_series("Actual", df_combined["Actual"])
    chart_data.add_series("Forecast", df_combined["Forecast"])
    chart_data.add_series("Upper", df_combined["Upper"])
    chart_data.add_series("Lower", df_combined["Lower"])

    x, y, cx, cy = Inches(0.5), Inches(1), Inches(9), Inches(6)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data).chart

    chart.chart_title.text_frame.text = "{0} forecast".format(i)

# %%
prs.save("report/tmc_target_meds/forecast_pmd_slides.pptx")
