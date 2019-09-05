from datetime import date
import glob
import numpy as np
import pandas as pd

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

filepaths = glob.glob("data/tidy/teresa/med-utilization_monthly_teresa_*.csv")
df = pd.concat(map(lambda x: pd.read_csv(x, index_col=0, parse_dates=True), filepaths), sort=False)
df.sort_index(inplace=True)
df['MONTH'] = pd.to_datetime('2018-' + (df.index-pd.DateOffset(months=6)).strftime('%m-%d')) + pd.DateOffset(months=6)
df['FY'] = 'FY' + (df.index + pd.DateOffset(months=6)).strftime('%y')

when = df.index[-1]
date_cut = date(day=1, month=7, year=(when+pd.DateOffset(months=6)).year) - pd.DateOffset(years=4)
df = df.loc[date_cut:]

df['MEDICATION'] = df['MEDICATION'].str.title()
df['MEDICATION'] = df['MEDICATION'].str.replace('Acetaminophen', 'Acetaminophen IV')
df['MEDICATION'] = df['MEDICATION'].str.replace('Levothyroxine', 'Levothyroxine IV')
df['MEDICATION'] = df['MEDICATION'].str.replace(' Human', '')
df['MEDICATION'] = df['MEDICATION'].str.replace('Prothrombin Complex', 'PCC 4-Factor')
df['MEDICATION'] = df['MEDICATION'].str.replace('Coagulation Factor Viia', 'Factor VIIa')
df['MEDICATION'] = df['MEDICATION'].str.replace('Antihemophilic Factor', 'Factor VIII')
df['MEDICATION'] = df['MEDICATION'].str.replace('Immune Globulin Intravenous And Subcut', 'IVIG')
df['MEDICATION'] = df['MEDICATION'].str.replace('Immune Globulin Intravenous', 'IVIG')

nurse_units = df['NURSE_UNIT'].unique()
meds = np.sort(df['MEDICATION'].unique())

df_pvt = pd.pivot_table(data=df, 
                        values='DOSES', 
                        index=['MEDICATION', 
                               'NURSE_UNIT', 
                               pd.Grouper(key='MONTH', freq='MS')], 
                        columns=pd.Grouper(key='FY'), 
                        aggfunc=np.sum, 
                        fill_value=0,
                        dropna=False)

prs = Presentation("doc/template.pptx")
# title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Medication Utilization in 7-Jones and Stroke"
data_end = df.index[-1].strftime('%B %Y')
subtitle.text = "Data through: " + data_end + "\nBrian Gulbis, PharmD, BCPS"

# blank slide layout
blank_slide_layout = prs.slide_layouts[6]
# chart size
x, y, cx, cy = Inches(0.5), Inches(1), Inches(9), Inches(6)

for i in nurse_units:
    # add section header slide
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    title = slide.shapes.title
    title.text = "Utilization in " + i

    for j in meds:
        df_j = df_pvt.loc[pd.IndexSlice[[j], [i], :]]
        
        if np.sum(df_j.sum()) == 0:
            continue
        
        # add new slide for graph
        slide = prs.slides.add_slide(blank_slide_layout)

        chart_data = CategoryChartData()
        chart_data.categories = df_j.index.get_level_values(2)

        for k in range(0, len(df_j.columns)):
            if k == 3:
                ser = df_j.iloc[:, k]
                idx = pd.to_datetime('2018-' + (when-pd.DateOffset(months=6)).strftime('%m-%d')) + pd.DateOffset(months=6)
                idx_next = idx + pd.DateOffset(months=1)
                keep = ser.loc[pd.IndexSlice[:, :, :idx]]
                chg = ser.loc[pd.IndexSlice[:, :, idx_next:]].replace({0: None})
                keep = keep.append(chg)
                chart_data.add_series(df_j.columns[k], keep)
            else:
                chart_data.add_series(df_j.columns[k], df_j.iloc[:, k])
        
        chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data).chart
        chart.chart_title.text_frame.text = "{0} utilization ({1})".format(j, i)

prs.save("report/teresa/utilization_slides.pptx")

# save data to Excel
with pd.ExcelWriter('report/teresa/utilization_data.xlsx') as writer:
    for i in nurse_units:
        df_excel = df.loc[df['NURSE_UNIT'] == i].copy()
        df_excel.reset_index(inplace=True)
        df_excel['MONTH'] = df_excel['EVENT_DATE'].dt.strftime('%Y-%m')
        df_excel = pd.pivot_table(data=df_excel,
                                  values='DOSES', 
                                  index='MEDICATION', 
                                  columns='MONTH',
                                  aggfunc=np.sum,
                                  fill_value=0)
        
        df_excel.to_excel(writer, sheet_name=i, freeze_panes=(1, 1))
