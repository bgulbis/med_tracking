import glob
import pandas as pd

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

def read_data(file):
    filepaths = glob.glob("../data/tidy/" + file + "/" + file + "_daily_doses*.csv")
    df = pd.concat(map(lambda x: pd.read_csv(x, index_col=0, parse_dates=True), filepaths), sort=False)
    df.sort_index(inplace=True)
    return df

def prep_df(df):
    df = df.resample("MS").sum()
    df["FY"] = "FY" + (df.index + pd.DateOffset(months=6)).strftime('%y')
    df["idx"] = pd.to_datetime("2019-" + (df.index + pd.DateOffset(months=-6)).strftime('%m-%d')) + pd.DateOffset(months=6)
    df_pvt = df['2016-07-01':].pivot(index="idx", columns="FY", values="DOSES").replace({pd.np.nan: None})

    return df_pvt

def add_utilization_slide(p, df, med):
    blank_slide_layout = p.slide_layouts[6]
    slide = p.slides.add_slide(blank_slide_layout)

    chart_data = CategoryChartData()
    chart_data.categories = df.index

    # for i in range(0, len(df.columns)):
    for i in range(len(df.columns) - 1, -1, -1):
        chart_data.add_series(df.columns[i], df.iloc[:, i])

    x, y, cx, cy = Inches(0.5), Inches(1), Inches(9), Inches(6)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data).chart

    if med == "ivig":
        med_title = med.upper()
    elif med == "acetaminophen-iv":
        med_title = "Acetaminophen IV"
    elif med == "bupivacaine-liposomal":
        med_title = "Bupivacaine (liposomal)"
    elif med == "levothyroxine-iv":
        med_title = "Levothyroxine IV"
    else:
        med_title = med.title()

    chart.chart_title.text_frame.text = med_title + " utilization"
    # chart.category_axis.axis_title.text_frame.text = "Month"
    chart.value_axis.has_title = True
    chart.value_axis.axis_title.text_frame.text = "Doses per month"

    return p

meds = ["acetaminophen-iv",
        "albumin",
        "sugammadex",
        "bupivacaine-liposomal",
        "isoproterenol",
        "nicardipine",
        "levothyroxine-iv",
        "ivig",
        "calcitonin",
        "pegfilgrastim",
        "eculizumab",
        "ceftaroline",
        "ceftazidime-avibactam",
        "ceftolozane-tazobactam",
        "daptomycin",
        "ertapenem",
        "meropenem-vaborbactam"]

prs = Presentation("../doc/template.pptx")
# title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Utilization of Target Medications"
data_end = read_data(meds[0]).index[-1].strftime('%B %Y')
subtitle.text = "Data through: " + data_end + "\nBrian Gulbis, PharmD, BCPS"

for i in meds:
    df = read_data(i)
    df = prep_df(df)
    add_utilization_slide(prs, df, i)

prs.save("../report/utilization/python_utilization_slides.pptx")
