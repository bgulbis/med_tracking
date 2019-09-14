from datetime import date
import glob
import numpy as np
import pandas as pd

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

df_days = pd.read_csv('data/tidy/patient_days.csv', index_col=0, parse_dates=True)
df_days.sort_index(inplace=True)

filepaths = glob.glob("data/tidy/brian/med-utilization_monthly_brian_*.csv")
df = pd.concat(map(lambda x: pd.read_csv(x, index_col=0, parse_dates=True), filepaths), sort=False)
df.sort_index(inplace=True)
df['MONTH'] = pd.to_datetime('2018-' + (df.index-pd.DateOffset(months=6)).strftime('%m-%d')) + pd.DateOffset(months=6)
df['FY'] = 'FY' + (df.index + pd.DateOffset(months=6)).strftime('%y')

when = df.index[-1]
date_cut = date(day=1, month=7, year=(when+pd.DateOffset(months=6)).year) - pd.DateOffset(years=4)
df = df.loc[date_cut:]

df['MEDICATION'] = df['MEDICATION'].str.title()
df['MEDICATION'] = df['MEDICATION'].str.replace('Acetaminophen', 'Acetaminophen IV')
df['MEDICATION'] = df['MEDICATION'].str.replace('Levothyroxine', 'Levothyroxine IV')
df['MEDICATION'] = df['MEDICATION'].str.replace(' Human', '')
df['MEDICATION'] = df['MEDICATION'].str.replace('Prothrombin Complex', 'PCC 4-Factor')
df['MEDICATION'] = df['MEDICATION'].str.replace('Coagulation Factor Viia', 'Factor VIIa')
df['MEDICATION'] = df['MEDICATION'].str.replace('Antihemophilic Factor', 'Factor VIII')
df['MEDICATION'] = df['MEDICATION'].str.replace('Immune Globulin Intravenous And Subcut', 'IVIG')
df['MEDICATION'] = df['MEDICATION'].str.replace('Immune Globulin Intravenous', 'IVIG')

nurse_units = df['NURSE_UNIT'].unique()
meds = np.sort(df['MEDICATION'].unique())

df['PATIENT_DAYS'] = df_days.loc[df_days['NURSE_UNIT'] == 'HH CVICU']['PATIENT_DAYS']
df['DOSES_PT_DAYS'] = df['DOSES']/df['PATIENT_DAYS']

df_pvt = pd.pivot_table(data=df, 
                        values='DOSES', 
                        index=['MEDICATION', 
                               'NURSE_UNIT', 
                               pd.Grouper(key='MONTH', freq='MS')], 
                        columns=pd.Grouper(key='FY'), 
                        aggfunc=np.sum, 
                        fill_value=0,
                        dropna=False)

prs = Presentation("doc/template.pptx")
# title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Medication Utilization in CVICU"
data_end = df.index[-1].strftime('%B %Y')
subtitle.text = "Data through: " + data_end + "\nBrian Gulbis, PharmD, BCPS"

# blank slide layout
blank_slide_layout = prs.slide_layouts[6]
# chart size
x, y, cx, cy = Inches(0.5), Inches(1), Inches(9), Inches(6)

for med in meds:
    df_med = df_pvt.loc[pd.IndexSlice[[med], ['HH CVICU'], :]]
    
    if np.sum(df_med.sum()) == 0:
        continue
    
    # add new slide for graph
    slide = prs.slides.add_slide(blank_slide_layout)

    chart_data = CategoryChartData()
    chart_data.categories = df_med.index.get_level_values(2)

    for i in range(0, len(df_med.columns)):
        if i == 3:
            ser = df_med.iloc[:, i]
            idx = pd.to_datetime('2018-' + (when-pd.DateOffset(months=6)).strftime('%m-%d')) + pd.DateOffset(months=6)
            idx_next = idx + pd.DateOffset(months=1)
            keep = ser.loc[pd.IndexSlice[:, :, :idx]]
            chg = ser.loc[pd.IndexSlice[:, :, idx_next:]].replace({0: None})
            keep = keep.append(chg)
            chart_data.add_series(df_med.columns[i], keep)
        else:
            chart_data.add_series(df_med.columns[i], df_med.iloc[:, i])

    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data).chart
    chart.chart_title.text_frame.text = "{0} doses per month".format(med)

    # add new slide for doses per patient days graph
    df_med_days = df.loc[df['MEDICATION'] == med].copy()

    slide = prs.slides.add_slide(blank_slide_layout)
    chart_data_days = CategoryChartData()
    chart_data_days.categories = df_med_days.index
    chart_data_days.add_series(df_med_days.columns[-1], df_med_days.iloc[:, -1])
    chart_days = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data_days).chart
    chart_days.chart_title.text_frame.text = "{0} doses per patient days".format(med)

prs.save("report/brian/utilization_slides.pptx")

# save data to Excel
with pd.ExcelWriter('report/brian/utilization_data.xlsx') as writer:
    for i in nurse_units:
        df_excel = df.loc[df['NURSE_UNIT'] == i].copy()
        df_excel.reset_index(inplace=True)
        df_excel['MONTH'] = df_excel['EVENT_DATE'].dt.strftime('%Y-%m')
        df_excel = pd.pivot_table(data=df_excel,
                                  values='DOSES', 
                                  columns='MEDICATION', 
                                  index='MONTH',
                                  aggfunc=np.sum,
                                  fill_value=0)
        
        df_excel.to_excel(writer, sheet_name='Doses', freeze_panes=(1, 1))

        df_excel = df.loc[df['NURSE_UNIT'] == i].copy()
        df_excel.reset_index(inplace=True)
        df_excel['MONTH'] = df_excel['EVENT_DATE'].dt.strftime('%Y-%m')
        df_excel = pd.pivot_table(data=df_excel,
                                  values='DOSES_PT_DAYS', 
                                  columns='MEDICATION', 
                                  index='MONTH',
                                  aggfunc=np.sum,
                                  fill_value=0)
        
        df_excel.to_excel(writer, sheet_name='Doses per Patient Days', freeze_panes=(1, 1))
