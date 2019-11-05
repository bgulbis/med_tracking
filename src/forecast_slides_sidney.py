from datetime import date
import glob
import numpy as np
import pandas as pd

from datetime import datetime
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from rpy2.robjects import pandas2ri
from rpy2.robjects.packages import importr
from rpy2.robjects.functions import SignatureTranslatedFunction

pandas2ri.activate()
forecast = importr('forecast')
forecast.auto_arima = SignatureTranslatedFunction(forecast.auto_arima, init_prm_translate = {'lambda_': 'lambda'})

filepaths = glob.glob("data/tidy/sidney/med-utilization_monthly_sidney_*.csv")
df = pd.concat(map(lambda x: pd.read_csv(x, index_col=0, parse_dates=True), filepaths), sort=False)
df.sort_index(inplace=True)

df['MEDICATION'] = df['MEDICATION'].str.title()
df['MEDICATION'] = df['MEDICATION'].str.replace('Acetaminophen', 'Acetaminophen IV')
df['MEDICATION'] = df['MEDICATION'].str.replace('Levothyroxine', 'Levothyroxine IV')
df['MEDICATION'] = df['MEDICATION'].str.replace(' Human', '')
df['MEDICATION'] = df['MEDICATION'].str.replace('Bupivacaine Liposome', 'Bupivacaine (liposomal)')
df['MEDICATION'] = df['MEDICATION'].str.replace('Immune Globulin Intravenous And Subcut', 'IVIG')
df['MEDICATION'] = df['MEDICATION'].str.replace('Immune Globulin Intravenous', 'IVIG')
df['MEDICATION'] = df['MEDICATION'].str.replace(' \\+ Sodium Chloride 0.9\\% Iv .*Ml', '', regex=True)

df['INPT'] = df['ENCOUNTER_TYPE'].isin(['Inpatient', 'Observation', 'Emergency'])
df.loc[(df['INPT'] == True) & (df['MEDICATION'] == "IVIG"), 'MEDICATION'] = 'IVIG (inpatient)'
df.loc[(df['INPT'] == False) & (df['MEDICATION'] == "IVIG"), 'MEDICATION'] = 'IVIG (outpatient)'

meds = np.sort(df['MEDICATION'].unique())

when = df.index[-1]
date_cut = date(day=1, month=7, year=(when+pd.DateOffset(months=6)).year) - pd.DateOffset(years=4)

prs = Presentation('doc/template.pptx')
# title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Forecast for Target Medications"
fcast_start = (df.index[-1] + pd.DateOffset(months=1)).strftime('%B %Y')
fcast_end = (df.index[-1] + pd.DateOffset(months=12)).strftime('%B %Y')
subtitle.text = "{0} to {1}\nBrian Gulbis, PharmD, BCPS".format(fcast_start, fcast_end)

for i in meds:
    df_actual = df.loc[df['MEDICATION'] == i].copy()

    if df_actual['DOSES'].sum() == 0:
        continue

    df_actual = df_actual.resample('MS').sum()
    df_actual.drop(['PATIENTS', 'INPT'], axis=1, inplace=True)
    df_actual.columns = ["Actual"]
    
    model = forecast.auto_arima(
        df_actual["Actual"],
        max_order=10,
        stepwise=False,
        approximation=False,
        lambda_="auto",
        biasadj=True
    )

    fcast = forecast.forecast(model, h=12)
    
    idx = pd.date_range(df_actual.index[-1] + pd.DateOffset(months=1), periods=12, freq='MS')
    yhat = pd.Series(fcast.rx2("mean"), name="Forecast", index=idx)
    lwr = pd.Series(fcast.rx2("lower")[0:12], name="Lower", index=idx)
    upr = pd.Series(fcast.rx2("upper")[0:12], name="Upper", index=idx)
    # lwr = pd.Series(fcast.rx2("lower")[:, :1].reshape(12, ), name="Lower", index=idx)
    # upr = pd.Series(fcast.rx2("upper")[:, :1].reshape(12, ), name="Upper", index=idx)
    df_fcast = pd.concat([yhat, lwr, upr], axis=1, sort=False)

    df_combined = pd.concat([df_actual, df_fcast], axis=1, sort=True)
    df_combined = df_combined.replace({pd.np.nan: None})
    df_combined = df_combined.loc[date_cut:]
        
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chart_data = CategoryChartData()
    chart_data.categories = df_combined.index

    chart_data.add_series("Actual", df_combined["Actual"])
    chart_data.add_series("Forecast", df_combined["Forecast"])
    chart_data.add_series("Upper", df_combined["Upper"])
    chart_data.add_series("Lower", df_combined["Lower"])

    x, y, cx, cy = Inches(0.5), Inches(1), Inches(9), Inches(6)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data).chart

    chart.chart_title.text_frame.text = "{0} forecast".format(i)

prs.save("report/sidney/forecast_slides.pptx")
