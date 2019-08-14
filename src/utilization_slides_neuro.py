from datetime import date
import glob
import numpy as np
import pandas as pd

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

def read_data(file):
    filepaths = glob.glob("../data/tidy/" + file + "/" + file + "_monthly_*.csv")
    df = pd.concat(map(lambda x: pd.read_csv(x, index_col=0, parse_dates=True), filepaths), sort=False)
    df.sort_index(inplace=True)
    return df

def prep_df(df, doses=True, encntr_type=None, facility=None, nurse_units=None):
    df['FY'] = 'FY' + (df.index + pd.DateOffset(months=6)).strftime('%y')
    df['MONTH'] = pd.to_datetime('2018-' + (df.index - pd.DateOffset(months=6)).strftime('%m-%d')) + pd.DateOffset(months=6)

    when = df.index[-1]
    date_cut = date(day=1, month=7, year=(when + pd.DateOffset(months=6)).year) - pd.DateOffset(years=4)
    df = df.loc[date_cut:]
    
    if encntr_type is not None:
        if type(encntr_type) is str:
                encntr_type = [encntr_type]
        df = df.loc[df['ENCOUNTER_TYPE'].isin(encntr_type)]
        
    if facility is not None:
        if type(facility) is str:
                facility = [facility]
        df = df.loc[df['FACILITY'].isin(facility)]
        
    if nurse_units is not None:
        if type(nurse_units) is str:
                nurse_units = [nurse_units]
        df = df.loc[df['NURSE_UNIT'].isin(nurse_units)]
    
    if doses == True:
        val = 'DOSES'
        df.drop('PATIENTS', axis=1, inplace=True)
    else:
        val = 'PATIENTS'
        df.drop('DOSES', axis=1, inplace=True)
        
    df = pd.pivot_table(df, values=val, index='MONTH', columns='FY', aggfunc=np.sum, fill_value=0, dropna=False)
# set the remaining months of current FY to None
#     df = df.replace({pd.np.nan: None})

    return df

def add_utilization_slide(p, df, med, title_post=""):
    blank_slide_layout = p.slide_layouts[6]
    slide = p.slides.add_slide(blank_slide_layout)

    chart_data = CategoryChartData()
    chart_data.categories = df.index

    # for i in range(0, len(df.columns)):
    for i in range(len(df.columns) - 1, -1, -1):
        chart_data.add_series(df.columns[i], df.iloc[:, i])

    x, y, cx, cy = Inches(0.5), Inches(1), Inches(9), Inches(6)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data).chart

    if med == "ivig":
        med_title = med.upper()
    elif med == "acetaminophen-iv":
        med_title = "Acetaminophen IV"
    elif med == "bupivacaine-liposomal":
        med_title = "Bupivacaine (liposomal)"
    elif med == "levothyroxine-iv":
        med_title = "Levothyroxine IV"
    else:
        med_title = med.title()

    chart.chart_title.text_frame.text = med_title + " utilization" + title_post
    # chart.category_axis.axis_title.text_frame.text = "Month"
    chart.value_axis.has_title = True
    chart.value_axis.axis_title.text_frame.text = "Doses per month"

    return p

meds = ["abciximab",
        "acetaminophen-iv",
        "albumin",
        "alteplase",
        "amphotericin",
        "ampicillin",
        "cefazolin",
        "cefepime",
        "ceftaroline",
        "ceftriaxone",
        "daptomycin",
        "dexmedetomidine",
        "eptifibatide",
        "ertapenem",
        "esmolol",
        "kcentra",
        "factor7",
        "factor8",
        "ivig",
        "levothyroxine-iv",
        "meropenem",
        "metronidazole",
        "micafungin",
        "nafcillin",
        "nicardipine",
        "nimodipine",
        "norepinephrine",
        "pentobarbital",
        "piperacillin-tazobactam",
        "propofol",
        "vancomycin",
        "vasopressin"]

meds = ['acetaminophen-iv']
nurse_units = ['HH 7J', 'HH STRK']

prs = Presentation("../doc/template.pptx")
# title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Medication Utilization in 7-Jones and Stroke"
data_end = read_data(meds[0]).index[-1].strftime('%B %Y')
subtitle.text = "Data through: " + data_end + "\nBrian Gulbis, PharmD, BCPS"

for i in meds:
    df = read_data(i)
    
    for j in nurse_units:
        df_j = prep_df(df, nurse_units=j)
        add_utilization_slide(prs, df_j, i, " ({0})".format(j))

prs.save("../report/utilization/utilization_slides_neuro.pptx")
