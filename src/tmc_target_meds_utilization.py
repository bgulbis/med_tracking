# %%
from datetime import date, datetime
import glob
import numpy as np
import pandas as pd

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

# %%
# import data
filepaths = glob.glob("data/raw/tmc_target_meds_*.csv")
df = pd.concat(map(lambda x: pd.read_csv(x, parse_dates=['DOSE_DATETIME']), filepaths), sort=False)

# %%
df['MEDICATION'] = df['MEDICATION'].str.title()
df['MEDICATION'] = df['MEDICATION'].str.replace('Acetaminophen', 'Acetaminophen IV')
df['MEDICATION'] = df['MEDICATION'].str.replace('Levothyroxine', 'Levothyroxine IV')
df['MEDICATION'] = df['MEDICATION'].str.replace(' Human', '')
df['MEDICATION'] = df['MEDICATION'].str.replace('Bupivacaine Liposome', 'Bupivacaine (liposomal)')
df['MEDICATION'] = df['MEDICATION'].str.replace('Immune Globulin Intravenous And Subcut', 'IVIG')
df['MEDICATION'] = df['MEDICATION'].str.replace('Immune Globulin Intravenous', 'IVIG')

df['INPT'] = df['ENCNTR_TYPE'].isin(['Inpatient', 'Observation', 'Emergency'])
df.loc[(df['INPT'] == True) & (df['MEDICATION'] == "IVIG"), 'MEDICATION'] = 'IVIG (inpatient)'
df.loc[(df['INPT'] == False) & (df['MEDICATION'] == "IVIG"), 'MEDICATION'] = 'IVIG (outpatient)'

df = df.sort_values(['MEDICATION', 'DOSE_DATETIME'])

meds = np.sort(df['MEDICATION'].unique())

# %%
df_doses = df.copy()
df_doses = df_doses.groupby('MEDICATION').resample('MS', on='DOSE_DATETIME').count()
df_doses = df_doses[['EVENT_ID']].rename(columns={'EVENT_ID': 'DOSES'})
df_doses = df_doses.reset_index(level='DOSE_DATETIME')
df_doses['MONTH'] = pd.to_datetime('2020-' + (df_doses['DOSE_DATETIME'] - pd.DateOffset(months=6)).dt.strftime('%m-%d')) + pd.DateOffset(months=6)
df_doses['FY'] = 'FY' + (df_doses['DOSE_DATETIME'] + pd.DateOffset(months=6)).dt.strftime('%y')

when = df['DOSE_DATETIME'].max()
date_cut = date(day=1, month=7, year=(when + pd.DateOffset(months=6)).year) - pd.DateOffset(years=4)
df_doses = df_doses[df_doses['DOSE_DATETIME'] >= date_cut]

# %%
df_pvt = pd.pivot_table(data=df_doses, 
                        values='DOSES', 
                        index=['MEDICATION', pd.Grouper(key='MONTH', freq='MS')], 
                        columns=pd.Grouper(key='FY'), 
                        aggfunc=np.sum, 
                        fill_value=0,
                        dropna=False)

# %%
# utilization slides
prs = Presentation("doc/template.pptx")

# title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Utilization of Target Medications"
data_end = when.strftime('%B %Y')
subtitle.text = "Data through: " + data_end + "\nBrian Gulbis, PharmD, BCPS"

# blank slide layout
blank_slide_layout = prs.slide_layouts[6]
# chart size
x, y, cx, cy = Inches(0.5), Inches(1), Inches(9), Inches(6)

# %%
for i in meds:
    df_i = df_pvt.loc[i]

    if np.sum(df_i.sum()) == 0:
        continue

    # add new slide for graph
    slide = prs.slides.add_slide(blank_slide_layout)

    chart_data = CategoryChartData()
    chart_data.categories = df_i.index

    for j in range(0, len(df_i.columns)):
        if j == 3:
            ser = df_i.iloc[:, j]
            idx = pd.to_datetime('2020-' + (when - pd.DateOffset(months=6)).strftime('%m-%d')) + pd.DateOffset(months=6)
            idx_next = idx + pd.DateOffset(months=1)
            keep = ser.loc[:idx]
            chg = ser.loc[idx_next:].replace({0: None})
            keep = keep.append(chg)
            chart_data.add_series(df_i.columns[j], keep)
        else:
            chart_data.add_series(df_i.columns[j], df_i.iloc[:, j])

    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data).chart
    chart.chart_title.text_frame.text = "{0} utilization".format(i)

# %%
df_sug = df.copy()
df_sug = df_sug[df_sug['MEDICATION'] == 'Sugammadex']

# %%
for i in ['MED_SERVICE', 'NURSE_UNIT', 'ENCNTR_TYPE']:
    for j in [0, 5]:
        m = datetime(day=1, month=(when - pd.DateOffset(months=j)).month, year=(when - pd.DateOffset(months=j)).year)
        df_sug_i = df_sug[df_sug['DOSE_DATETIME'] >= m]
        df_sug_i = df_sug_i.groupby(i).count()
        df_sug_i = df_sug_i[['EVENT_ID']].rename(columns={'EVENT_ID': 'DOSES'})
        df_sug_i = df_sug_i.sort_values(by='DOSES')
        df_sug_i = df_sug_i.tail(20)

        # add new slide for graph
        slide = prs.slides.add_slide(blank_slide_layout)

        chart_data = CategoryChartData()
        chart_data.categories = df_sug_i.index
        chart_data.add_series('Doses', df_sug_i['DOSES'])
        chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data).chart
        
        if df_sug_i.shape[0] < 20:
            t = df_sug_i.shape[0]
        else:
            t = 20

        svc = {'MED_SERVICE': 'service lines', 'NURSE_UNIT': 'nurse units', 'ENCNTR_TYPE': 'encounter types'}
        dt = {0: "in " + m.strftime('%b %Y'), 5: "since " + m.strftime('%b %Y')}
        chart.chart_title.text_frame.text = f"Top {t} {svc[i]} utilizing sugammadex {dt[j]}"

# %%
prs.save("report/tmc_target_meds/utilization_slides.pptx")