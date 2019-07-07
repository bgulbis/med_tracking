import pandas as pd
from pptx import Presentation 
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import (XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION, 
                            XL_MARKER_STYLE, XL_TICK_MARK)
from pptx.enum.dml import MSO_LINE_DASH_STYLE, MSO_THEME_COLOR
from pptx.util import Inches, Pt

def add_forecast_slide(p, df, med):
    df = df.replace({pd.np.nan: None})
    blank_slide_layout = p.slide_layouts[6]
    slide = p.slides.add_slide(blank_slide_layout)

    chart_data = CategoryChartData()
    chart_data.categories = df.index
    chart_data.add_series("Actual", df["Actual"])
    chart_data.add_series("Forecast", df["Forecast"])

    x, y, cx, cy = Inches(1), Inches(1), Inches(8), Inches(6)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data).chart
    chart = format_graph(chart, med)
    return p

def format_graph(chart, med, font_nm = None, ax_bright = 0.35):
    category_axis = format_axis_title(chart.category_axis, "Month", font_nm, ax_bright)
    category_axis = format_axis_tick_labels(category_axis, font_nm, ax_bright, "mmm yy")
    category_axis = format_axis(category_axis, ax_bright)
#     category_axis.maximum_scale = len(df_combined) + 1
#     df_combined.asfreq("MS").index[-1] + pd.DateOffset(months=1)

    value_axis = format_axis_title(chart.value_axis, "Number of doses", font_nm, ax_bright)
    value_axis = format_axis_tick_labels(value_axis, font_nm, ax_bright, "#,##0")
    value_axis = format_axis(value_axis, ax_bright)

    if med == "ivig":
        med_title = med.upper()
    elif med == "acetaminophen":
        med_title = med.title() + " IV "
    else:
        med_title = med.title()
    
    chart.has_title = True
    chart.chart_title.text_frame.text = med_title + " forecast"
    chart.chart_title.text_frame.paragraphs[0].font.name = font_nm
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(24)
    chart.chart_title.text_frame.paragraphs[0].font.bold = False
    chart.chart_title.text_frame.paragraphs[0].font.color.theme_color = MSO_THEME_COLOR.TEXT_1
    chart.chart_title.text_frame.paragraphs[0].font.color.brightness = 0.25

    chart.has_legend = False
    # chart.legend.include_in_layout = False
    # chart.legend.position = XL_LEGEND_POSITION.TOP

    # format Actual line
    col_actual = MSO_THEME_COLOR.ACCENT_1
    chart.series[0].format.line.width = Pt(3.5)
    chart.series[0].format.line.color.theme_color = col_actual

    i = len(chart.series[0].points) - 13
    chart = format_marker(chart, 0, i, font_nm, MSO_THEME_COLOR.BACKGROUND_1, col_actual)

    # format Forecast line
    brght_forecast = 0.4

    chart.series[1].format.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    chart.series[1].format.line.width = Pt(2.5)
    chart.series[1].format.line.color.theme_color = col_actual
    chart.series[1].format.line.color.brightness = brght_forecast

    j = len(chart.series[1].points) - 1
    chart = format_marker(chart, 1, j, font_nm, col_actual, col_actual, brght_forecast)
    
    return chart

def format_axis(axis, ax_bright):
    axis.format.line.color.theme_color = MSO_THEME_COLOR.TEXT_1
    axis.format.line.color.brightness = ax_bright
    axis.has_major_gridlines = False
    axis.has_minor_gridlines = False
    return axis

def format_axis_title(axis, title_text, font_nm, ax_bright, has_title = True):
    axis.has_title = has_title
    axis.axis_title.text_frame.text = title_text
    axis.axis_title.text_frame.paragraphs[0].font.name = font_nm
    axis.axis_title.text_frame.paragraphs[0].font.size = Pt(18)
    axis.axis_title.text_frame.paragraphs[0].font.bold = False
    axis.axis_title.text_frame.paragraphs[0].font.color.theme_color = MSO_THEME_COLOR.TEXT_1
    axis.axis_title.text_frame.paragraphs[0].font.color.brightness = ax_bright
    return axis

def format_axis_tick_labels(axis, font_nm, ax_bright, num_fmt):
    axis.tick_labels.font.name = font_nm
    axis.tick_labels.font.size = Pt(16)
    axis.tick_labels.font_bold = False
    axis.tick_labels.font.color.theme_color = MSO_THEME_COLOR.TEXT_1
    axis.tick_labels.font.color.brightness = ax_bright
    axis.tick_labels.number_format = num_fmt
    return axis

def format_marker(chart, i, j, font_nm, col_mark, col_line, brght=0):
    chart.series[i].data_labels.number_format = "#,##0"
#     chart.series[i].data_labels.show_value = True
    chart.series[i].data_labels.font.name = font_nm
    chart.series[i].data_labels.font.size = Pt(18)
    chart.series[i].data_labels.font.bold = False
    chart.series[i].data_labels.position = XL_DATA_LABEL_POSITION.BELOW
    chart.series[i].points[j].marker.style = XL_MARKER_STYLE.CIRCLE
    chart.series[i].points[j].marker.size = 7
    chart.series[i].points[j].marker.format.fill.solid()
    chart.series[i].points[j].marker.format.fill.fore_color.theme_color = col_mark
    chart.series[i].points[j].marker.format.fill.fore_color.brightness = brght
    chart.series[i].points[j].marker.format.line.width = Pt(2.5)
    chart.series[i].points[j].marker.format.line.color.theme_color = col_line
    chart.series[i].points[j].marker.format.line.color.brightness = brght
    return chart
