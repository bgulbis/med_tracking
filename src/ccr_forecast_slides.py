import glob
import numpy as np
import pandas as pd

from datetime import datetime
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from rpy2.robjects import pandas2ri
from rpy2.robjects.packages import importr
from rpy2.robjects.functions import SignatureTranslatedFunction

pandas2ri.activate()
forecast = importr('forecast')
forecast.auto_arima = SignatureTranslatedFunction(forecast.auto_arima, init_prm_translate = {'lambda_': 'lambda'})

def read_data(file):
    filepaths = glob.glob("../data/tidy/" + file + "/" + file + "_daily_doses*.csv")
    df = pd.concat(map(lambda x: pd.read_csv(x, index_col=0, parse_dates=True), filepaths), sort=False)
    df.sort_index(inplace=True)
    return df

def make_forecast(df, h=12):
    f_m = forecast.auto_arima(
        df["DOSES"],
        max_order=10,
        stepwise=False,
        approximation=False,
        lambda_="auto",
        biasadj=True
    )

    fc_m = forecast.forecast(f_m, h=h)
    return fc_m

def make_forecast_df(med, h=12, resamp="MS", sep="all"):
    inpt = ["Inpatient", "Observation"]

    df = read_data(med)
    
    if sep == "inpt":
        df = df.loc[df["ENCOUNTER_TYPE"].isin(inpt)]
    elif sep == "outpt":
        df = df.loc[~df["ENCOUNTER_TYPE"].isin(inpt)]
    
    # df_encntr = df.copy()
    # 
    # if "ENCOUNTER_TYPE" in df_encntr.columns:
    #     df_encntr["Inpatient"] = df["ENCOUNTER_TYPE"].str.contains("inpatient|observation", case=False)
    #     df_pvt = df_encntr.pivot_table(values="DOSES", index="EVENT_DATE", columns="Inpatient", aggfunc=np.sum)
    #     df_pvt = df_pvt.resample("MS").sum()
    #     df_pvt.columns = ["Outpatient", "Inpatient"]

    df = df.resample(resamp).sum()
    fc = make_forecast(df, h)
    idx_m = pd.date_range(df.index[-1] + pd.DateOffset(months=1), periods=h, freq=resamp)
    yhat = pd.Series(fc.rx2("mean"), name="Forecast", index=idx_m)
    lwr = pd.Series(fc.rx2("lower")[:, :1].reshape(12, ), name="Lower", index=idx_m)
    upr = pd.Series(fc.rx2("upper")[:, :1].reshape(12, ), name="Upper", index=idx_m)
    df_fc = pd.concat([yhat, lwr, upr], axis=1, sort=False)
    df.columns = ["Actual"]

    # if "ENCOUNTER_TYPE" in df_encntr.columns:
    #     x = [df, df_fc, df_pvt]
    # else:
    x = [df, df_fc]

    df_combined = pd.concat(x, axis=1).replace({pd.np.nan: None})

    return df_combined

def add_forecast_slide(p, df, med, title_post=""):
    blank_slide_layout = p.slide_layouts[6]
    slide = p.slides.add_slide(blank_slide_layout)

    # num_fmt = "#,##0"
    chart_data = CategoryChartData()
    chart_data.categories = df.index

    chart_data.add_series("Actual", df["Actual"])
    chart_data.add_series("Forecast", df["Forecast"])
    chart_data.add_series("Upper", df["Upper"])
    chart_data.add_series("Lower", df["Lower"])

    x, y, cx, cy = Inches(0.5), Inches(1), Inches(9), Inches(6)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data).chart

    if med == "ivig":
        med_title = med.upper()
    elif med == "acetaminophen-iv":
        med_title = "Acetaminophen IV"
    elif med == "bupivacaine-liposomal":
        med_title = "Bupivacaine (liposomal)"
    elif med == "levothyroxine-iv":
        med_title = "Levothyroxine IV"
    else:
        med_title = med.title()

    chart.chart_title.text_frame.text = med_title + " forecast" + title_post
    chart.value_axis.axis_title.text_frame.text = "Doses per month"

    return p

def make_slides(meds, n=12):
    prs = Presentation()
    # title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Forecast for Target Medications"
    fcast_start = datetime.now().strftime('%B %Y')
    fcast_end = (datetime.now() + pd.DateOffset(months=11)).strftime('%B %Y')
    subtitle.text = fcast_start + " to " + fcast_end + "\nBrian Gulbis, PharmD, BCPS"

    for i in meds:
        df = make_forecast_df(i)
        add_forecast_slide(prs, df, i)

    return prs

meds = ["acetaminophen-iv",
        "albumin",
        "sugammadex",
        "bupivacaine-liposomal",
        "isoproterenol",
        "nicardipine",
        "levothyroxine-iv",
        "ivig",
        "calcitonin",
        "pegfilgrastim",
        "eculizumab",
        "ceftaroline",
        "ceftazidime-avibactam",
        "ceftolozane-tazobactam",
        "daptomycin",
        "ertapenem",
        "meropenem-vaborbactam"]

prs = Presentation("../doc/template.pptx")
# title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Forecast for Target Medications"
fcast_start = datetime.now().strftime('%B %Y')
fcast_end = (datetime.now() + pd.DateOffset(months=11)).strftime('%B %Y')
subtitle.text = fcast_start + " to " + fcast_end + "\nBrian Gulbis, PharmD, BCPS"

for i in meds:
    if i == "ivig":
        df = make_forecast_df(i, sep="inpt")
        add_forecast_slide(prs, df, i, " (inpatient)")

        df = make_forecast_df(i, sep="outpt")
        add_forecast_slide(prs, df, i, " (outpatient)")
        
    else:
        df = make_forecast_df(i)
        add_forecast_slide(prs, df, i)

prs.save("../report/utilization/python_forecast_slides.pptx")
